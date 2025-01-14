from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Optional
import mlx.core as mx
from mlx_lm import load, generate
import json
import uvicorn
import pandas as pd
import csv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

app = FastAPI()

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Load the model (Mistral 7B by default)
try:
    model, tokenizer = load("mlx-community/Mistral-7B-v0.1-hf-4bit-mlx")
except Exception as e:
    raise RuntimeError(f"Error loading the model: {e}")

def process_file(file: UploadFile):
    try:
        content = file.file.read()
        if file.filename.endswith('.csv'):
            decoded_content = content.decode('utf-8')
            reader = csv.reader(decoded_content.splitlines())
            data = [row for row in reader]
            return f"CSV File Content:\n{data}"
        elif file.filename.endswith('.json'):
            json_data = json.loads(content)
            return f"JSON File Content:\n{json.dumps(json_data, indent=2)}"
        elif file.filename.endswith(('.xls', '.xlsx')):
            df = pd.read_excel(file.file)
            return f"Excel File Content:\n{df.to_string(index=False)}"
        elif file.filename.endswith('.txt'):
            decoded_content = content.decode('utf-8')
            return f"Text File Content:\n{decoded_content}"
        elif file.filename.endswith('.pdf'):
            file.file.seek(0)  # Reset file pointer
            reader = PdfReader(file.file)
            pdf_text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
            return f"PDF File Content:\n{pdf_text}"
        elif file.filename.endswith('.docx'):
            file.file.seek(0)  # Reset file pointer
            doc = Document(file.file)
            doc_text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            return f"Word File Content:\n{doc_text}"
        elif file.filename.endswith('.pptx'):
            file.file.seek(0)  # Reset file pointer
            presentation = Presentation(file.file)
            slides_text = []
            for slide in presentation.slides:
                slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                slides_text.append(slide_text)
            return f"PowerPoint File Content:\n{'\n'.join(slides_text)}"
        else:
            return f"Unsupported file format: {file.filename}"
    except Exception as e:
        return f"Error processing file {file.filename}: {e}"

async def process_request(message: str, files: Optional[List[UploadFile]] = None):
    try:
        # Process files if present
        file_contents = []
        if files:
            for file in files:
                file_content = process_file(file)
                file_contents.append(file_content)

        # Prepare prompt
        prompt = f"User: {message}"
        if file_contents:
            prompt += "\n\nAttached files:\n" + "\n".join(file_contents)
        prompt += "\n\nAssistant:"

        # Generate response
        generated_text = generate(model, tokenizer, prompt=prompt, max_tokens=500)

        if isinstance(generated_text, str):
            return {"text": generated_text}
        else:
            raise ValueError(f"Unexpected output from generate: {type(generated_text)}")

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating response: {str(e)}")

@app.post("/chat")
async def chat(
    message: str = Form(...),
    files: List[UploadFile] = File(None)
):
    print(f"Received message: {message}")  # Debug log
    print(f"Received files: {files}")  # Debug log

    if not message or message.strip() == "":
        raise HTTPException(status_code=400, detail="Message is required")

    response = await process_request(message, files)
    return response

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
